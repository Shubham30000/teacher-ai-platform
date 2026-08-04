"""
PPTX parser built on ``python-pptx``.

Each slide becomes one top-level :class:`Section` (heading = slide
title, or "Slide N" if the layout has no title placeholder). Text
placeholders/body boxes become paragraphs or list items depending on
paragraph indentation level; tables and pictures on the slide are
preserved on that section; speaker notes are appended as a labeled
paragraph so they are not silently dropped.
"""
from __future__ import annotations

import logging
import re
from pathlib import Path

from app.core.constants import ElementType, SupportedFileType
from app.core.exceptions import ParsingError
from app.document_intelligence.models import (
    ContentBlock,
    DocumentMetadata,
    HeadingLevel,
    ImageElement,
    Section,
    StructuredDocument,
    TableCell,
    TableElement,
    UrlElement,
)
from app.parsers.base import BaseParser

logger = logging.getLogger(__name__)

_URL_RE = re.compile(r"https?://[^\s)\]}'\"<>]+")


class PptxParser(BaseParser):
    format_name = "pptx"
    supports_ocr = False

    def parse(self, file_path: Path) -> StructuredDocument:
        self._base_check(file_path)
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise ParsingError(
                "python-pptx is not installed. Install it with `pip install python-pptx`."
            ) from exc

        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:  # noqa: BLE001
            raise ParsingError(f"Failed to open PPTX file: {exc}") from exc

        sections: list[Section] = []
        out_dir = file_path.parent / f"{file_path.stem}_assets"
        image_index = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_text = None
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()

            section = Section(
                heading=title_text or f"Slide {slide_number}",
                heading_level=HeadingLevel.H1,
                page_number=slide_number,
            )

            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
                        if not text:
                            continue
                        block_type = (
                            ElementType.LIST_ITEM if paragraph.level > 0 else ElementType.PARAGRAPH
                        )
                        section.blocks.append(
                            ContentBlock(
                                type=block_type,
                                text=text,
                                page_number=slide_number,
                                list_level=paragraph.level,
                            )
                        )
                        for url in _URL_RE.findall(text):
                            section.urls.append(UrlElement(url=url, page_number=slide_number))

                if shape.has_table:
                    table = shape.table
                    n_rows = len(table.rows)
                    n_cols = len(table.columns)
                    cells = [
                        TableCell(
                            text=table.cell(r, c).text,
                            row=r,
                            col=c,
                            is_header=(r == 0),
                        )
                        for r in range(n_rows)
                        for c in range(n_cols)
                    ]
                    section.tables.append(
                        TableElement(
                            n_rows=n_rows, n_cols=n_cols, cells=cells, page_number=slide_number
                        )
                    )

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_index += 1
                    stored_path = None
                    try:
                        out_dir.mkdir(parents=True, exist_ok=True)
                        ext = shape.image.ext or "png"
                        out_path = out_dir / f"slide{slide_number}_image_{image_index}.{ext}"
                        out_path.write_bytes(shape.image.blob)
                        stored_path = str(out_path)
                    except Exception:  # noqa: BLE001 - image extraction is best-effort
                        logger.warning(
                            "Could not extract image %d on slide %d of '%s'",
                            image_index,
                            slide_number,
                            file_path.name,
                        )
                    section.images.append(
                        ImageElement(page_number=slide_number, stored_path=stored_path)
                    )

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    section.blocks.append(
                        ContentBlock(
                            type=ElementType.PARAGRAPH,
                            text=f"[Speaker notes] {notes_text}",
                            page_number=slide_number,
                        )
                    )

            sections.append(section)

        metadata = DocumentMetadata(
            source_filename=file_path.name,
            file_type=SupportedFileType.PPTX,
            file_size_bytes=file_path.stat().st_size,
            slide_count=len(sections),
            title=(presentation.core_properties.title or None),
            author=(presentation.core_properties.author or None),
        )
        doc = StructuredDocument(metadata=metadata, sections=sections)
        logger.info(
            "Parsed PPTX '%s': %d slides, %d tables, %d images",
            file_path.name,
            len(sections),
            doc.total_table_count(),
            doc.total_image_count(),
        )
        return doc
